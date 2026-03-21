import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Criar DataFrames de exemplo
data1 = {'Nome': ['Alice', 'Bob', 'Charlie'],
         'Idade': [25, 20, 22]}
df1 = pd.DataFrame(data1)

data2 = {'Produto': ['Maçã', 'Banana', 'Laranja'],
         'Quantidade': [10, 15, 8]}
df2 = pd.DataFrame(data2)

# Criar apresentação PowerPoint
prs = Presentation()

# Função para adicionar DataFrame como tabela ao slide
def add_dataframe_to_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Slide de conteúdo de título e conteúdo
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    shape = slide.shapes.add_table(rows=df.shape[0] + 1, cols=df.shape[1], left=left, top=top, width=width, height=height).table

    # Adicionar cabeçalho
    for col_num, col_name in enumerate(df.columns):
        cell = shape.cell(0, col_num)
        cell.text = col_name

    # Adicionar dados
    for row_num, row in df.iterrows():
        for col_num, value in enumerate(row):
            cell = shape.cell(row_num + 1, col_num)
            cell.text = str(value)

# Adicionar DataFrames aos slides
add_dataframe_to_slide(prs, df1)
add_dataframe_to_slide(prs, df2)

# Salvar apresentação
output_path = 'dataframes_presentation.pptx'
prs.save(output_path)
print(f'Apresentação salva em {output_path}')
